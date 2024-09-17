from http.client import HTTPException
import traceback
from supabase_utils.pptx_utils import read_pptx_from_supabase
from fastapi import FastAPI, Request, WebSocket, WebSocketDisconnect, UploadFile, File
from fastapi.responses import Response, JSONResponse
from fastapi.staticfiles import StaticFiles
from unstructured.partition.pdf import partition_pdf
from pydantic import BaseModel
import json
import os
import io
import logging
from typing import List
from reach_core.utils.websocket_manager import WebSocketManager
from reach_core.master.prompts import component_injection, generate_report_prompt
from fastapi.middleware.cors import CORSMiddleware
from pptx.util import Inches, Pt
from openai import OpenAI
import subprocess
import tempfile
import os


# from reach_core.utils.unstructured_functions import *
# from reach_core.utils.hubspot_functions import *
# from reach_core.utils.whisper_functions import *
from reach_core.utils.mp3_from_mp4 import mp4_to_mp3
# from utils import write_md_to_pdf
from datetime import datetime

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ResearchRequest(BaseModel):
    task: str
    report_type: str
    agent: str
    sources: List[str] = []

class SalesforceCredentials(BaseModel):
    username: str
    consumer_key: str
    private_key_path: str

class HubspotCredentials(BaseModel):
    access_token: str

def make_serializable(data):
    if isinstance(data, dict):
        return {key: make_serializable(value) for key, value in data.items()}
    elif isinstance(data, list):
        return [make_serializable(item) for item in data]
    elif isinstance(data, datetime):
        return data.isoformat()
    elif hasattr(data, '__dict__'):
        return {key: make_serializable(value) for key, value in data.__dict__.items()}
    else:
        return data

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://themagi.systems", "https://www.themagi.systems", "wss://themagi.systems"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class TaskRequest(BaseModel):
    task: str
    report_type: str
    sources: List[str]

manager = WebSocketManager()


# Dynamic directory for outputs once first research is run
@app.on_event("startup")
def startup_event():
    if not os.path.isdir("outputs"):
        os.makedirs("outputs")
    app.mount("/outputs", StaticFiles(directory="outputs"), name="outputs")

# @app.get("/")
# async def read_root(request: Request):
#     return templates.TemplateResponse('index.html', {"request": request, "report": None})

@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    await manager.connect(websocket)
    logger.info("WebSocket connected")
    # logger.info(f"Received WebSocket data: {data}")
    try:
        while True:
            data = await websocket.receive_text()
            try:
                json_data = json.loads(data)
                print(f"raw data {json_data}")
                task = json_data.get("task")
                report_type = json_data.get("report_type")
                sources = json_data.get("sources", [])
                edits = json_data.get("edits")  #None if not present
                cadence = json_data.get("cadence")
                print(f"ws edits {edits}")
                if task and report_type:
                    await manager.start_streaming(task, report_type, sources, websocket, cadence, edits)
                    await websocket.send_json({"type": "complete"})
                else:
                    await websocket.send_json({"type": "error", "message": "Missing required parameters"})
            except json.JSONDecodeError:
                logging.error(f"Invalid JSON received: {data}")
                await websocket.send_json({"type": "error", "message": "Invalid JSON format"})
            except Exception as e:
                logging.error(f"Error processing request: {str(e)}")
                await websocket.send_json({"type": "error", "message": "Internal server error"})
    except WebSocketDisconnect:
        await manager.disconnect(websocket)

client = OpenAI()

class SlideContent(BaseModel):
    title: str
    content: List[str]

class PresentationStructure(BaseModel):
    title: str
    slides: List[SlideContent]

class PowerPointRequest(BaseModel):
    prompt: str
    filePath: str
    favoriteTheme: str
    signedUrl: str

@app.post("/generate-powerpoint")
async def generate_powerpoint(request: PowerPointRequest):
    try:
        print(f"Received request with prompt: {request.prompt}")
        print(f"File path: {request.filePath}")
        print(f"Favorite theme: {request.favoriteTheme}")
        print(f"Signed URL: {request.signedUrl}")
        
        prs = read_pptx_from_supabase(request.filePath, request.signedUrl) 

        user_prompt = f"""
            Create a presentation summarizing the following content. Be sure to include the links in a References slide: {request.prompt}.
            Use headings to drive the narrative for the presentation and slide titles. Distill the core content down to create succinct slide content.
        """  

        completion = client.chat.completions.create(
            model="gpt-4o-2024-08-06",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates PowerPoint presentations. Generate a presentation structure based on the user's request by calling the create_presentation function."},
                {"role": "user", "content": user_prompt}
            ],
            functions=[{
                "name": "create_presentation",
                "description": "Create a PowerPoint presentation structure",
                "parameters": PresentationStructure.schema()
            }],
            function_call={"name": "create_presentation"}
        )
        print("OpenAI API response received")

        function_call = completion.choices[0].message.function_call
        if not function_call or function_call.name != "create_presentation":
            raise ValueError("Unexpected response from OpenAI API")

        presentation_data = json.loads(function_call.arguments)
        print(f"Presentation data: {presentation_data}")

        trim_count = len(prs.slides)
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        title = title_slide.shapes.title
        subtitle = next((ph for ph in title_slide.placeholders if ph.placeholder_format.idx != 0), None)
        
        if title:
            title.text = presentation_data["title"]
        if subtitle:
            subtitle.text = "Generated Presentation"
        
        content_slide_layout = prs.slide_layouts[1]
        for slide_data in presentation_data["slides"]:
            new_slide = prs.slides.add_slide(content_slide_layout)
            
            title = new_slide.shapes.title
            content_placeholder = next((ph for ph in new_slide.placeholders if ph.placeholder_format.idx != 0), None)
            
            if title:
                title.text = slide_data["title"]
            
            if content_placeholder:
                tf = content_placeholder.text_frame
            else:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                txBox = new_slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame

            tf.clear()
            for item in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0

        for _ in range(trim_count):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        print("Presentation generated successfully")

        return Response(
            content=ppt_bytes.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=generated_presentation.pptx"}
        )
    except Exception as e:
        print(f"Error in generate_powerpoint: {str(e)}")
        print(traceback.format_exc())
        return JSONResponse(status_code=500, content={"detail": str(e)})

class AnalysisPoint(BaseModel):
    key_point: str
    important_language: List[str]

class AmbiguousClause(BaseModel):
    clause: str
    ambiguous_language: List[str]

class DocumentAnalysis(BaseModel):
    key_points: List[AnalysisPoint]
    ambiguous_clauses: List[AmbiguousClause]

@app.post("/process-pdf")
async def process_pdf(file: UploadFile = File(...)):
    with open(file.filename, "wb") as buffer:
        buffer.write(await file.read())
    
    elements = partition_pdf(file.filename)
    text_content = "\n".join([str(el) for el in elements])
    print(f'debugging: {text_content}')
    
    os.remove(file.filename)
    
    response = client.chat.completions.create(
        model="gpt-4o-2024-08-06",
        messages=[
            {"role": "system", "content": "You are an expert legal assistant that analyzes documents and extracts key information. Format the content by calling the analyze_document function."},
            {"role": "user", "content": f"Analyze the following document and provide key points with important language, and ambiguous clauses with their language. \n\nDocument content:\n{text_content}"}
        ],
        functions=[{
            "name": "analyze_document",
            "description": "Analyze the document and provide structured output",
            "parameters": DocumentAnalysis.schema()
        }],
        function_call={"name": "analyze_document"}
    )
    
    function_call = response.choices[0].message.function_call
    if not function_call or function_call.name != "analyze_document":
        raise ValueError("Unexpected response from OpenAI API")

    analysis = json.loads(function_call.arguments)
    
    return analysis

class ChartRequest(BaseModel):
    tableId: str
    tableContent: str
    chartConfig: dict
    previousError: str = None

@app.post("/create-chart")
async def create_chart(request: ChartRequest):
    component_code = component_injection()
    config_str = "\n".join([f"{k}: {v}" for k, v in request.chartConfig.items()])
    
    base_prompt = f"""
    You are an expert D3.js developer.
    Create only the necessary D3.js code to generate a chart based on the following data and configuration:
    Data:
    {request.tableContent}
    Chart Configuration:
    {config_str}
    Follow these guidelines:
    - Use the specified color scheme. If not specified, lean into earth toned colors like stone-400, stone-500, zinc-600, and zinc-500.
    - Include a legend if showLegend is true.
    - Add a title if showTitle is true, using the specified titleFontSize.
    - Include axis labels if showAxisLabels is true.
    - If showChartLabels is true, add labels to the chart elements.
    - Avoid the use of `translate`, `width`, and `selectAll`.
    - If the x-axis contains time indices, ensure time increases to the right.
    - Do not redeclare the variable `svg` or any other variables if they have already been declared in the environment.
    - Assume that a D3.js environment is already available and that an `svg` element has been appended to the DOM.
    - Do not include any HTML tags, <script> tags, or references to external libraries.
    - Avoid pie charts entirely.

    Provide only the D3.js code, without any explanation or additional text.

    The D3.js code will be executed in the following component in renderChart: {component_code}.
    The D3 code's compatibility with the provided component is mission critical.

    Be creative! Use a chart type that best visualizes the data and tells the most rich data story.
    """

    if request.previousError:
        prompt = f"{base_prompt}\n\nThe previous attempt resulted in an error: {request.previousError}. Please correct the D3.js code to ensure it is valid and can be rendered correctly."
    else:
        prompt = base_prompt

    try:
        completion = client.chat.completions.create(
            model="gpt-4o-2024-08-06",
            messages=[
                {"role": "system", "content": "You are an expert D3.js developer."},
                {"role": "user", "content": prompt}
            ],
        )
        d3_code = completion.choices[0].message.content
        return {"d3_code": d3_code}
    except Exception as e:
        print(f"Error in create-chart: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error generating chart code")

class CondenseRequest(BaseModel):
    task: str
    accumulatedOutput: str
    
@app.post("/condense-findings")
async def condense_findings(request: CondenseRequest):
    user_prompt = generate_report_prompt(
        question=request.task,
        context=request.accumulatedOutput
    )

    prompt = f"""
    You are a seasoned analyst. \n
    Your primary goal is to compose comprehensive, astute, impartial, and methodically arranged reports of the provided research.\n
    Aggregate key thematic points, quantitative findings, and tables from the provided research.
    """

    try:
        completion = client.chat.completions.create(
            model="gpt-4o-2024-08-06",
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": user_prompt}
            ],
        )
        
        condensed_report = completion.choices[0].message.content

        return {"condensed_report": condensed_report}
    except Exception as e:
        print(f"Error in condense-findings: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500)
    
class DiagramRequest(BaseModel):
    content: str
    diagramType: str
    previousError: str = None

@app.post("/generate-diagram")
async def generate_diagram(request: DiagramRequest):
    component_code = component_injection()

    try:
        diagram_examples = {
            "flowchart": """
            flowchart TD
                A[Start] --> B{Is it?}
                B -->|Yes| C[OK]
                C --> D[Rethink]
                D --> B
                B ---->|No| E[End]

            """,
            "quadrantChart": """
            quadrantChart
                title Reach and engagement of campaigns
                x-axis Low Reach --> High Reach
                y-axis Low Engagement --> High Engagement
                quadrant-1 We should expand
                quadrant-2 Need to promote
                quadrant-3 Re-evaluate
                quadrant-4 May be improved
                Campaign A: [0.3, 0.6]
                Campaign B: [0.45, 0.23]
                Campaign C: [0.57, 0.69]
                Campaign D: [0.78, 0.34]
                Campaign E: [0.40, 0.34]
                Campaign F: [0.35, 0.78]

            """,
            "c4Context": """
                C4Context
                    title System Context diagram for Internet Banking System
                    Enterprise_Boundary(b0, "BankBoundary0") {
                        Person(customerA, "Banking Customer A", "A customer of the bank, with personal bank accounts.")
                        Person(customerB, "Banking Customer B")
                        Person_Ext(customerC, "Banking Customer C", "desc")

                        Person(customerD, "Banking Customer D", "A customer of the bank, <br/> with personal bank accounts.")

                        System(SystemAA, "Internet Banking System", "Allows customers to view information about their bank accounts, and make payments.")

                        Enterprise_Boundary(b1, "BankBoundary") {

                        SystemDb_Ext(SystemE, "Mainframe Banking System", "Stores all of the core banking information about customers, accounts, transactions, etc.")

                        System_Boundary(b2, "BankBoundary2") {
                            System(SystemA, "Banking System A")
                            System(SystemB, "Banking System B", "A system of the bank, with personal bank accounts. next line.")
                        }

                        System_Ext(SystemC, "E-mail system", "The internal Microsoft Exchange e-mail system.")
                        SystemDb(SystemD, "Banking System D Database", "A system of the bank, with personal bank accounts.")

                        Boundary(b3, "BankBoundary3", "boundary") {
                            SystemQueue(SystemF, "Banking System F Queue", "A system of the bank.")
                            SystemQueue_Ext(SystemG, "Banking System G Queue", "A system of the bank, with personal bank accounts.")
                        }
                        }
                    }

                    BiRel(customerA, SystemAA, "Uses")
                    BiRel(SystemAA, SystemE, "Uses")
                    Rel(SystemAA, SystemC, "Sends e-mails", "SMTP")
                    Rel(SystemC, customerA, "Sends e-mails to")

                    UpdateElementStyle(customerA, $fontColor="red", $bgColor="grey", $borderColor="red")
                    UpdateRelStyle(customerA, SystemAA, $textColor="blue", $lineColor="blue", $offsetX="5")
                    UpdateRelStyle(SystemAA, SystemE, $textColor="blue", $lineColor="blue", $offsetY="-10")
                    UpdateRelStyle(SystemAA, SystemC, $textColor="blue", $lineColor="blue", $offsetY="-40", $offsetX="-50")
                    UpdateRelStyle(SystemC, customerA, $textColor="red", $lineColor="red", $offsetX="-50", $offsetY="20")

                    UpdateLayoutConfig($c4ShapeInRow="3", $c4BoundaryInRow="1")
            """,
            "timeline": """
            timeline
                title History of AI Development
                section 1950s
                    1956 : Dartmouth Conference
                        : Birth of AI as a field
                section 1960s-1970s
                    1965 : ELIZA chatbot
                    1972 : PROLOG language
                section 1980s-1990s
                    1997 : IBM's Deep Blue beats Kasparov
                section 2000s-2010s
                    2011 : IBM Watson wins Jeopardy!
                    2014 : Google acquires DeepMind
                section 2020s
                    2022 : ChatGPT released
            """,
            "mindmap": """
              mindmap
                root((mindmap))
                    Origins
                    Long history
                    ::icon(fa fa-book)
                    Popularisation
                        British popular psychology author Tony Buzan
                    Research
                    On effectiveness<br/>and features
                    On Automatic creation
                        Uses
                            Creative techniques
                            Strategic planning
                            Argument mapping
                    Tools
                    Pen and paper
                    Mermaid
            """
        }

        base_prompt = f"""
        You are an expert in creating Mermaid.js diagrams. Based on the following content, create a Mermaid.js {request.diagramType} diagram that best represents the information:

        {request.content}

        Here's a syntax example for the {request.diagramType}:

        ```mermaid
        {diagram_examples[request.diagramType]}
        ```

        Important notes:
        - The first line of your code MUST be the diagram type: {request.diagramType}
        - Adhere closely to the syntax provided in the example.
        - Ensure the diagram accurately represents the provided content.
        - Do not nest parentheses.
        - Do not use abbreviations in parenthesis. Ex: 'Complete Blood Count (CBC)' is invalid syntax.
        - NEVER include a title denoted with '%% Title:'. Ex: '%% Title: Diagnostic Process for Thoracic Pain'
        - Do NOT use underscores in entity names or labels. Underscores are protected syntax in Mermaid C4Context diagrams and should not be used as text. Use spaces instead.

        The code will be rendered in the following component with renderDiagram, its compatibility is mission critical: {component_code}
        Provide only the Mermaid.js code, without any explanation or additional text.
        """

        if request.previousError:
            prompt = f"{base_prompt}\n\nThe previous attempt resulted in an error. Please correct the Mermaid.js code to ensure it is valid and can be rendered correctly. Remember to start with the diagram type: {request.diagramType}"
        else:
            prompt = base_prompt

        client = OpenAI()
        completion = client.chat.completions.create(
            model="gpt-4o-2024-08-06",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates Mermaid.js diagrams."},
                {"role": "user", "content": prompt}
            ]
        )

        mermaid_code = completion.choices[0].message.content.strip()
        
        if not mermaid_code.startswith(request.diagramType):
            mermaid_code = f"{request.diagramType}\n{mermaid_code}"

        return {"mermaid_code": mermaid_code}
    except Exception as e:
        print(f"Error in generate-diagram: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# @app.post("/setEnvironmentVariables")
# async def set_environment_variables(credentials: SalesforceCredentials):

#     parsed_contents = await process_salesforce(
#         username=credentials.username,
#         consumer_key=credentials.consumer_key,
#         private_key=credentials.private_key_path
#     )
    
#     parsed_contents = make_serializable(parsed_contents)

#     salesforce_dir = "salesforce"
#     if not os.path.isdir(salesforce_dir):
#         os.makedirs(salesforce_dir)

#     parsed_uploads_path = f"{salesforce_dir}/parsed_app.json"
#     if not os.path.exists(parsed_uploads_path):
#         async with aiofiles.open(parsed_uploads_path, "w") as new_file:
#             await new_file.write("[]")

#     async with aiofiles.open(parsed_uploads_path, "r+") as parsed_uploads_file:
#         existing_content = await parsed_uploads_file.read()
#         existing_data = json.loads(existing_content) if existing_content else []
#         existing_data.extend(parsed_contents)
#         await parsed_uploads_file.seek(0)
#         await parsed_uploads_file.write(json.dumps(existing_data))
#         await parsed_uploads_file.truncate()


#     return {"message": "Environment variables set successfully and app processed"}

# @app.post("/setEnvironmentVariables")
# async def set_environment_variables(credentials: HubspotCredentials):
    
#     parsed_contents = await process_hubspot_crm_objects(credentials=credentials.access_token)

#     parsed_contents = make_serializable(parsed_contents)

#     hubspot_dir = "hubspot"
#     if not os.path.isdir(hubspot_dir):
#         os.makedirs(hubspot_dir)

#     parsed_uploads_path = f"{hubspot_dir}/parsed_app.json"
#     if not os.path.exists(parsed_uploads_path):
#         async with aiofiles.open(parsed_uploads_path, "w") as new_file:
#             await new_file.write("[]")

#     async with aiofiles.open(parsed_uploads_path, "r+") as parsed_uploads_file:
#         existing_content = await parsed_uploads_file.read()
#         existing_data = json.loads(existing_content) if existing_content else []
#         existing_data.extend(parsed_contents)
#         await parsed_uploads_file.seek(0)
#         await parsed_uploads_file.write(json.dumps(existing_data))
#         await parsed_uploads_file.truncate()


#     return {"message": "Environment variables set successfully and app processed"}

# @app.post("/upload")
# async def upload_files(files: List[UploadFile] = File(...), task: str = Form(...)):
#     upload_dir = "uploads"
#     if not os.path.isdir(upload_dir):
#         os.makedirs(upload_dir)

#     async def handle_file(file):
#         file_location = f"{upload_dir}/{file.filename}"
#         with open(file_location, "wb+") as file_object:
#             content = await file.read()
#             file_object.write(content)
        
#         if file.filename.endswith(".mp4"):
#             mp3_paths = await mp4_to_mp3(file_location)
#             os.remove(file_location)
#             file_location = mp3_paths[0]
        
#         return file_location

#     uploaded_files_info = []
#     for file in files:
#         file_location = await handle_file(file)
#         uploaded_files_info.append({"filename": file.filename, "location": file_location})
#         print(f"Uploaded: {file.filename} to {file_location}")

#     parsed_contents = []
#     for uploaded_file_info in uploaded_files_info:
#         parsed_audio = await parse_text_from_audio()
#         parsed_documents = await process_unstructured()
#         parsed_content = parsed_audio + parsed_documents
#         parsed_contents.extend(parsed_content)

#     parsed_uploads_path = f"{upload_dir}/parsed_uploads.json"
#     if not os.path.exists(parsed_uploads_path):
#         async with aiofiles.open(parsed_uploads_path, "w") as new_file:
#             await new_file.write("[]")

#     async with aiofiles.open(parsed_uploads_path, "r+") as parsed_uploads_file:
#         existing_content = await parsed_uploads_file.read()
#         existing_data = json.loads(existing_content) if existing_content else []
#         existing_data += parsed_contents
#         await parsed_uploads_file.seek(0)
#         await parsed_uploads_file.write(json.dumps(existing_data))
#         await parsed_uploads_file.truncate()

#     return {"info": f"Files uploaded successfully", "task": task, "parsed_info": f"Data written to {parsed_uploads_path}"}
