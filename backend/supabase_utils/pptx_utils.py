import os
from dotenv import load_dotenv
from supabase import create_client, Client
import requests
from PyPDF2 import PdfReader
from pptx import Presentation
import io
from pptx.slide import Slide
from pptx.slide import SlideLayout
from pptx.presentation import Presentation as PresentationType
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.dml import MSO_THEME_COLOR


load_dotenv()

supabase: Client = create_client(
    supabase_url=os.getenv("NEXT_PUBLIC_SUPABASE_URL"),
    supabase_key=os.getenv("NEXT_PUBLIC_SUPABASE_ANON_KEY")
    )

def read_pptx_from_supabase(file_path: str, signed_url: str) -> Presentation:
    try:
        print(f"Attempting to download file from: {signed_url}")
        file_response = requests.get(signed_url)
        file_response.raise_for_status()
        ppt = Presentation(io.BytesIO(file_response.content))
        
        print("Template structure:")
        for i, slide in enumerate(ppt.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                print(f"  Shape: {shape.name}, Type: {shape.shape_type}")
        
        return ppt
    except Exception as e:
        print(f"Error reading PowerPoint from Supabase: {str(e)}")
        raise

def extract_text_from_pdf(pdf_content: bytes) -> str:
    pdf_reader = PdfReader(io.BytesIO(pdf_content))
    text_content = ""
    for page_num, page in enumerate(pdf_reader.pages, 1):
        text_content += f"\n\n--- Page {page_num} ---\n\n"
        text_content += page.extract_text()
    return text_content